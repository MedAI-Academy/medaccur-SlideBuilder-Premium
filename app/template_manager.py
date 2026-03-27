"""Template Manager — loads, caches, and validates PPTX designer templates."""
from __future__ import annotations

import os
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

TEMPLATE_DIR = Path(__file__).parent.parent / "templates"

# ══════════════════════════════════════════════════════════
# SHAPE NAME SPEC — every template must have these shapes
# ══════════════════════════════════════════════════════════
# Each section type requires certain named shapes in its template.
# Designers set shape names via PowerPoint > Selection Pane.

REQUIRED_SHAPES = {
    "_title": ["title_drug", "title_subtitle", "title_year", "title_footer"],
    "_divider": ["divider_number", "divider_title", "divider_accent_bar"],
    "executive_summary": ["header_title", "card_1_title", "card_1_body", "card_2_title", "card_2_body", "card_3_title", "card_3_body"],
    "prevalence_kpi": ["header_title", "kpi_1_value", "kpi_1_label", "kpi_2_value", "kpi_2_label", "kpi_3_value", "kpi_3_label"],
    "treatment_algo": ["header_title", "table_area"],
    "guidelines": ["header_title", "table_area"],
    "unmet_needs": ["header_title", "card_1_title", "card_1_body"],
    "moa": ["header_title", "drug_class_label", "step_area"],
    "pivotal_table": ["header_title", "kpi_mpfs", "kpi_hr", "kpi_orr", "kpi_mos", "chart_area"],
    "swot": ["header_title", "s_title", "s_body", "w_title", "w_body", "o_title", "o_body", "t_title", "t_body"],
    "imperatives": ["header_title", "roof_title", "pillar_1_title", "pillar_1_body", "base_bar"],
    "competitor_table": ["header_title", "chart_area", "table_area"],
    "narrative": ["header_title", "primary_message", "talking_points"],
    "tactics": ["header_title", "table_area"],
    "timeline": ["header_title", "road_area"],
    "summary": ["header_title", "message_area"],
}


def list_available_themes() -> list[str]:
    """Return list of theme IDs that have template directories."""
    if not TEMPLATE_DIR.exists():
        return []
    return [d.name for d in TEMPLATE_DIR.iterdir() if d.is_dir() and not d.name.startswith(".")]


def get_template_path(theme: str, section_type: str) -> Path | None:
    """Get path to a specific section template, or None if not found."""
    path = TEMPLATE_DIR / theme / f"{section_type}.pptx"
    if path.exists():
        return path
    return None


def get_full_template_path(theme: str) -> Path | None:
    """Get path to the full 20-slide template for a theme, or None if not found."""
    path = TEMPLATE_DIR / theme / "_full.pptx"
    if path.exists():
        return path
    return None


def load_template(theme: str, section_type: str) -> Presentation | None:
    """Load a template PPTX for a given theme and section type."""
    path = get_template_path(theme, section_type)
    if path is None:
        return None
    return Presentation(str(path))


def find_shape_by_name(slide, name: str):
    """Find a shape on a slide by its name (set via Selection Pane in PowerPoint)."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_shape_text(shape, text: str, font_size: int = None, font_color: str = None, bold: bool = None):
    """Replace all text in a shape, preserving formatting where possible."""
    if shape is None or not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    # Clear existing paragraphs by setting first paragraph text and removing rest
    if tf.paragraphs:
        p = tf.paragraphs[0]
        if p.runs:
            run = p.runs[0]
            run.text = text
            if font_size:
                run.font.size = Pt(font_size)
            if font_color:
                run.font.color.rgb = RGBColor.from_string(font_color)
            if bold is not None:
                run.font.bold = bold
        else:
            p.text = text
        # Remove extra paragraphs
        while len(tf.paragraphs) > 1:
            p_elem = tf.paragraphs[-1]._p
            p_elem.getparent().remove(p_elem)


def validate_template(theme: str, section_type: str) -> dict:
    """Validate that a template has all required named shapes."""
    tpl = load_template(theme, section_type)
    if tpl is None:
        return {
            "valid": False,
            "template_id": f"{theme}/{section_type}",
            "missing_shapes": REQUIRED_SHAPES.get(section_type, []),
            "found_shapes": [],
            "message": f"Template file not found: {theme}/{section_type}.pptx",
        }

    required = REQUIRED_SHAPES.get(section_type, [])
    if not required:
        return {
            "valid": True,
            "template_id": f"{theme}/{section_type}",
            "missing_shapes": [],
            "found_shapes": [],
            "message": f"No shape requirements defined for {section_type}",
        }

    # Check first slide for named shapes
    if len(tpl.slides) == 0:
        return {
            "valid": False,
            "template_id": f"{theme}/{section_type}",
            "missing_shapes": required,
            "found_shapes": [],
            "message": "Template has no slides",
        }

    slide = tpl.slides[0]
    found = [shape.name for shape in slide.shapes]
    missing = [name for name in required if name not in found]

    return {
        "valid": len(missing) == 0,
        "template_id": f"{theme}/{section_type}",
        "missing_shapes": missing,
        "found_shapes": [n for n in found if n in required],
        "message": "OK" if not missing else f"Missing {len(missing)} required shapes",
    }
