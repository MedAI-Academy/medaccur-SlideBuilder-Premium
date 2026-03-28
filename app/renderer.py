"""Core PPTX Renderer — builds premium-quality slides with python-pptx.

Rendering modes:
  1. TEMPLATE MODE: Load designer .pptx, find shapes by name, fill data.
  2. PROGRAMMATIC MODE (default): Build slides from scratch with python-pptx.
     Higher quality than PptxGenJS because we control XML directly.

Each section_type has a dedicated render function.
"""
import io
import os
import re
import shutil
import tempfile
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import chart_renderer
from . import template_manager

import copy as copy_module
from lxml import etree


# ══════════════════════════════════════════════════════════
# TEMPLATE SLIDE CLONING — copy designer slide into target pres
# This is the KEY difference from the failed Java approach:
# We find shapes by .name (set in Selection Pane), NOT by text content.
# No XML run splitting problem because we replace the entire text frame.
# ══════════════════════════════════════════════════════════

def _clone_template_slide(pres: Presentation, template_pptx_path: str):
    """
    Clone the first slide from a template PPTX into the target presentation.
    Returns the new slide, or None if cloning fails.
    
    Copies shapes, background, AND image/media relationships so embedded
    pictures and backgrounds survive the clone.
    """
    try:
        tpl = Presentation(template_pptx_path)
        if not tpl.slides or len(tpl.slides) == 0:
            return None
        
        src_slide = tpl.slides[0]
        src_part = src_slide.part
        
        # Add blank slide to target
        blank_layout = pres.slide_layouts[6]  # Blank layout
        new_slide = pres.slides.add_slide(blank_layout)
        dest_part = new_slide.part
        
        # ── Step 1: Copy image/media relationships ──
        rId_map = {}
        for rId, rel in list(src_part.rels.items()):
            try:
                if 'image' in rel.reltype or 'media' in rel.reltype:
                    new_rId = dest_part.relate_to(rel.target_part, rel.reltype)
                    rId_map[rId] = new_rId
            except Exception:
                pass
        
        # Also copy images from source slide layout and master into target
        try:
            for ancestor in [src_slide.slide_layout, src_slide.slide_layout.slide_master]:
                if ancestor is None:
                    continue
                for rId, rel in list(ancestor.part.rels.items()):
                    try:
                        if 'image' in rel.reltype:
                            target_ancestor = new_slide.slide_layout if ancestor == src_slide.slide_layout else new_slide.slide_layout.slide_master
                            if target_ancestor:
                                target_ancestor.part.relate_to(rel.target_part, rel.reltype)
                    except Exception:
                        pass
        except Exception:
            pass
        
        # ── Step 2: Remove default shapes from blank slide ──
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        
        # ── Step 3: Copy each shape, updating rId references ──
        for shape in src_slide.shapes:
            el = copy_module.deepcopy(shape._element)
            # Update any rId references in the copied element
            if rId_map:
                for sub_el in el.iter():
                    for attr_name in list(sub_el.attrib.keys()):
                        val = sub_el.attrib[attr_name]
                        if val in rId_map:
                            sub_el.attrib[attr_name] = rId_map[val]
            new_slide.shapes._spTree.append(el)
        
        # ── Step 4: Copy background ──
        try:
            PML_NS = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
            src_cSld = src_slide._element.find(f'{PML_NS}cSld')
            dest_cSld = new_slide._element.find(f'{PML_NS}cSld')
            if src_cSld is not None and dest_cSld is not None:
                src_bg = src_cSld.find(f'{PML_NS}bg')
                if src_bg is not None:
                    new_bg = copy_module.deepcopy(src_bg)
                    # Update rId references in background
                    if rId_map:
                        for sub_el in new_bg.iter():
                            for attr_name in list(sub_el.attrib.keys()):
                                val = sub_el.attrib[attr_name]
                                if val in rId_map:
                                    sub_el.attrib[attr_name] = rId_map[val]
                    existing_bg = dest_cSld.find(f'{PML_NS}bg')
                    if existing_bg is not None:
                        dest_cSld.remove(existing_bg)
                    dest_cSld.insert(0, new_bg)
        except Exception:
            pass
        
        return new_slide
    except Exception as e:
        print(f"Template clone failed: {e}")
        import traceback
        traceback.print_exc()
        return None


def _find_shape(slide, name: str):
    """Find a shape on a slide by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _find_shapes_by_name(slide, name: str) -> list:
    """Find ALL shapes with a given name (handles duplicates like 'Oval 40')."""
    return [s for s in slide.shapes if s.name == name]


def _get_text_color_deep(shape):
    """
    Extract the DESIGNER'S intended text color from a shape's XML.
    Scans rPr (run properties) and defRPr (default run properties)
    for the first explicit color, resolving schemeClr via SLIDEWORKS_THEME.
    
    This is critical because:
    - Slideworks 2025 theme defines bg2 = lt2 = #FFFFFF
    - Many shapes inherit white text via schemeClr="bg2"
    - When we replace runs, the schemeClr reference is lost
    - So we must extract and re-apply the ORIGINAL color explicitly
    """
    if not hasattr(shape, 'text_frame'):
        return None
    A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    root = etree.fromstring(etree.tostring(shape._element))
    
    # Check rPr (explicit run properties) first
    for rPr in root.iter(f'{A}rPr'):
        sf = rPr.find(f'{A}solidFill')
        if sf is not None:
            sr = sf.find(f'{A}srgbClr')
            if sr is not None:
                return sr.get('val')
            sc = sf.find(f'{A}schemeClr')
            if sc is not None:
                return SLIDEWORKS_THEME.get(sc.get('val'), '000000')
    
    # Check defRPr (default run properties)
    for defRPr in root.iter(f'{A}defRPr'):
        sf = defRPr.find(f'{A}solidFill')
        if sf is not None:
            sr = sf.find(f'{A}srgbClr')
            if sr is not None:
                return sr.get('val')
            sc = sf.find(f'{A}schemeClr')
            if sc is not None:
                return SLIDEWORKS_THEME.get(sc.get('val'), '000000')
    
    return None


def _set_text(shape, text: str, force_color: str = None):
    """
    Replace ALL text in a shape, preserving the DESIGNER'S original text color.
    
    Key insight: The Slideworks template uses schemeClr references (bg2, tx1, etc.)
    for text colors. When python-pptx replaces runs, these references are lost.
    We extract the original color from XML BEFORE clearing, then re-apply it
    explicitly as srgbClr on every new run.
    
    Args:
        shape: PowerPoint shape with text_frame
        text: New text content (supports \n for multi-line)
        force_color: Optional hex color to override (e.g. "FFFFFF" for white)
    """
    if shape is None or not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    
    A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    
    # Get designer's intended color BEFORE clearing text
    color_hex = force_color or _get_text_color_deep(shape) or "2D3748"
    
    # Save formatting from first run of first paragraph
    first_p = tf.paragraphs[0]
    saved_font = {}
    if first_p.runs:
        r = first_p.runs[0]
        saved_font = {
            'size': r.font.size,
            'bold': r.font.bold,
            'italic': r.font.italic,
            'name': r.font.name,
        }
    saved_alignment = first_p.alignment
    
    # Clear all paragraphs except first
    p_elements = list(tf._txBody.iterchildren(f'{A_NS}p'))
    for p_el in p_elements[1:]:
        tf._txBody.remove(p_el)
    
    # Handle multi-line text
    lines = str(text).split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            for r_el in list(p._p.iterchildren(f'{A_NS}r')):
                p._p.remove(r_el)
        else:
            p = tf.add_paragraph()
        
        p.alignment = saved_alignment
        run = p.add_run()
        run.text = line
        
        # Apply explicit color (never rely on inheritance)
        run.font.color.rgb = RGBColor.from_string(color_hex)
        
        # Apply saved formatting
        if saved_font.get('size'):
            run.font.size = saved_font['size']
        if saved_font.get('bold') is not None:
            run.font.bold = saved_font['bold']
        if saved_font.get('italic') is not None:
            run.font.italic = saved_font['italic']
        if saved_font.get('name'):
            run.font.name = saved_font['name']


# ══════════════════════════════════════════════════════════
# TEMPLATE FILL FUNCTIONS — populate cloned slides with data
# Each function knows which shape names to look for and what data to put in.
# Falls back gracefully if shapes are missing.
# ══════════════════════════════════════════════════════════

def _fill_title(slide, meta: dict, ordered_count: int, ref_count: int):
    """Fill the title slide (slide 0)."""
    _set_text(_find_shape(slide, "title_drug"), meta.get("drug", "Medical Affairs Plan"))
    scope = meta.get("country") or meta.get("region") or "Global"
    _set_text(_find_shape(slide, "title_subtitle"), f"{scope} Medical Affairs Plan — {meta.get('year', '')}")
    _set_text(_find_shape(slide, "title_year"), meta.get("indication", ""))


def _fill_divider(slide, num: int, title: str):
    """Fill a divider slide."""
    _set_text(_find_shape(slide, "divider_number"), str(num).zfill(2))
    _set_text(_find_shape(slide, "divider_title"), title)


def _fill_executive_summary(slide, d: dict, meta: dict):
    _set_text(_find_shape(slide, "header_title"),
              f"Executive Summary — {meta.get('country','Global')} {meta.get('year','')}")
    rows = d.get("rows", [])
    body = "\n\n".join(
        f"{r.get('icon','')} {r.get('topic','')}\n{r.get('summary','')}"
        for r in rows[:8]
    )
    _set_text(_find_shape(slide, "body_text"), body or "No executive summary data.")


def _fill_disease_intro(slide, d: dict, meta: dict):
    _set_text(_find_shape(slide, "header_title"),
              f"Disease Introduction — {meta.get('indication', '')}")
    slides_data = d.get("slides", [{}])
    sl_data = slides_data[0] if slides_data else {}
    items = sl_data.get("items", [])
    if items:
        body = "\n\n".join(f"{_str(it.get('label',''))}\n{_str(it.get('text',''))}" for it in items[:8])
    else:
        lines = [f"{k}: {_str(v)}" for k, v in d.items() if k != "section_type" and not isinstance(v, (dict, list))]
        body = "\n\n".join(lines)
    _set_text(_find_shape(slide, "body_text"), body or "No disease introduction data.")


def _fill_prevalence_kpi(slide, d: dict, meta: dict):
    """Slide 5: 3 label shapes below circles + takeaway bar. Don't write in circle shapes (icons block them)."""
    _set_text(_find_shape(slide, "header_title"),
              f"Prevalence & Epidemiology — {meta.get('country','Global')}")
    kpis = d.get("kpis", [])
    # Labels below circles: combine label + value
    label_shapes = ["Google Shape;606;p17", "Google Shape;607;p17", "Google Shape;608;p17"]
    for i, k in enumerate(kpis[:3]):
        if i < len(label_shapes):
            _set_text(_find_shape(slide, label_shapes[i]),
                      f"{_str(k.get('label',''))}\n{_str(k.get('value',''))}")
    # Takeaway bar
    ctx = d.get("context_table", [])
    if ctx:
        summary = " | ".join(f"{_str(c.get('metric',''))}: {_str(c.get('value',''))}" for c in ctx[:3])
    else:
        summary = f"{len(kpis)} epidemiology metrics collected"
    _set_text(_find_shape(slide, "Rectangle 28"), summary)
    # Source
    sources = [_str(k.get('source','')) for k in kpis[:3] if k.get('source')]
    _set_text(_find_shape(slide, "Google Shape;507;p14"),
              f"Source: {', '.join(sources)[:80]}" if sources else "")


def _fill_moa(slide, d: dict, meta: dict):
    """Slide 10: step_1-4, Chevron 14, Rectangle 34-44 descriptions."""
    _set_text(_find_shape(slide, "header_title"),
              f"Mechanism of Action — {meta.get('drug','')}")
    _set_text(_find_shape(slide, "Subtitle 2"),
              f"{_str(d.get('drug_class',''))} · Target: {_str(d.get('target',''))}")
    steps = d.get("pathway_steps", [])
    title_shapes = ["step_1", "step_2", "Chevron 14", "step_3", "step_4"]
    desc_shapes = ["Rectangle 34", "Rectangle 36", "Rectangle 43", "Rectangle 44", "Rectangle 35"]
    for i, s in enumerate(steps[:5]):
        if i < len(title_shapes):
            _set_text(_find_shape(slide, title_shapes[i]), _str(s.get("title", f"Step {i+1}")))
        if i < len(desc_shapes):
            _set_text(_find_shape(slide, desc_shapes[i]), _str(s.get("description", "")))
    _set_text(_find_shape(slide, "Google Shape;507;p14"),
              f"Source: {_str(d.get('source', d.get('key_reference','')))[:80]}")


def _fill_pivotal_table(slide, d: dict, meta: dict):
    """Slide 4: KPIs in TextBoxes. Remove placeholder Chart 16."""
    _set_text(_find_shape(slide, "header_title"), f"Pivotal Studies — {meta.get('drug','')}")
    trials = d.get("trials", [])
    active = [t for t in trials if t.get("include_in_map") is not False]
    # Remove template placeholder chart (meaningless bar chart)
    chart_shape = _find_shape(slide, "Chart 16")
    if chart_shape:
        try:
            chart_shape._element.getparent().remove(chart_shape._element)
        except Exception:
            pass
    if active:
        tr = active[0]
        eff = tr.get("efficacy", {})
        _set_text(_find_shape(slide, "Subtitle 28"),
                  f"{_str(tr.get('name',''))} — {_str(tr.get('phase',''))} — {_str(tr.get('design',''))}")
        _set_text(_find_shape(slide, "TextBox 2"),
                  f"mPFS: {_str(eff.get('mpfs_drug','N/A'))} vs {_str(eff.get('mpfs_control','N/A'))} · HR {_str(eff.get('hr_pfs','N/A'))} · ORR {_str(eff.get('orr_drug','N/A'))}")
        # Clear overlapping TextBoxes
        _set_text(_find_shape(slide, "TextBox 10"), "")
        _set_text(_find_shape(slide, "TextBox 7"), "")
        _set_text(_find_shape(slide, "TextBox 8"),
                  f"N={_str(tr.get('n_total','N/A'))} · mOS: {_str(eff.get('mos_drug','NR'))}")
        saf = tr.get("safety", {})
        details = []
        if saf.get("grade34_heme"):
            details.append("G3-4 Heme: " + ", ".join(_str(s) for s in saf["grade34_heme"][:3]))
        if saf.get("grade34_nonheme"):
            details.append("G3-4 Non-heme: " + ", ".join(_str(s) for s in saf["grade34_nonheme"][:3]))
        if saf.get("discontinuation_rate"):
            details.append(f"D/C: {_str(saf['discontinuation_rate'])}")
        _set_text(_find_shape(slide, "TextBox 17"), "\n".join(details))
        _set_text(_find_shape(slide, "Google Shape;507;p14"), f"Source: {_str(tr.get('source',''))[:80]}")


def _fill_differentiators(slide, d: dict, meta: dict):
    """Slide 6: card_1-4 + clear 'Icon' text from Oval 12-15."""
    _set_text(_find_shape(slide, "header_title"), f"Key Differentiators — {meta.get('drug','')}")
    diffs = d.get("differentiators", [])
    for i, df in enumerate(diffs[:4]):
        shape = _find_shape(slide, f"card_{i+1}")
        if shape:
            _set_text(shape, f"{_str(df.get('title',''))}\n\n{_str(df.get('detail', df.get('evidence','')))}")
    # Clear "Icon" text from Ovals (decorative circles with picture overlay)
    for n in ["Oval 12", "Oval 13", "Oval 14", "Oval 15"]:
        _set_text(_find_shape(slide, n), "")


def _fill_areas_interest(slide, d: dict, meta: dict):
    """Slide 7: 5x (Text Placeholder 6 = area name, Rectangle 89 = description)."""
    _set_text(_find_shape(slide, "header_title"),
              f"Areas of Interest (ISR) — {meta.get('drug','')}")
    areas = d.get("areas", [])
    area_names = _find_shapes_by_name(slide, "Text Placeholder 6")
    area_descs = _find_shapes_by_name(slide, "Rectangle 89")
    for i in range(min(5, len(area_names))):
        if i < len(areas):
            area = areas[i]
            _set_text(area_names[i], _str(area.get("area", f"Area {i+1}")))
            interests = area.get("interests", [])
            desc = " · ".join(_str(it)[:60] for it in interests[:4])
            if i < len(area_descs):
                _set_text(area_descs[i], desc)
        else:
            _set_text(area_names[i], "")
            if i < len(area_descs):
                _set_text(area_descs[i], "")


def _fill_iep(slide, d: dict, meta: dict):
    """Slide 8: Fill existing 8x5 table_area."""
    _set_text(_find_shape(slide, "header_title"),
              f"Integrated Evidence Plan — {meta.get('drug','')} {meta.get('year','')}")
    table_shape = _find_shape(slide, "table_area")
    if table_shape is None:
        return
    try:
        tbl = table_shape.table
    except Exception:
        return
    gaps = d.get("gaps", [])
    for ri, g in enumerate(gaps[:tbl.rows.__len__() - 1]):
        data = [_str(g.get("gap",""))[:80], _str(g.get("activity",""))[:80],
                _str(g.get("responsible",""))[:40], _str(g.get("status",""))[:30],
                _str(g.get("source",""))[:50]]
        for ci in range(min(len(data), tbl.columns.__len__())):
            try:
                cell = tbl.cell(ri + 1, ci)
                cell.text = data[ci]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)
            except Exception:
                pass
    # Clear unused rows
    for ri in range(len(gaps) + 1, tbl.rows.__len__()):
        for ci in range(tbl.columns.__len__()):
            try:
                tbl.cell(ri, ci).text = ""
            except Exception:
                pass


def _fill_guidelines(slide, d: dict, meta: dict):
    """Slide 9: Fill existing 6x5 table_area."""
    _set_text(_find_shape(slide, "header_title"), f"Guidelines — {meta.get('drug','')}")
    table_shape = _find_shape(slide, "table_area")
    if table_shape is None:
        return
    try:
        tbl = table_shape.table
    except Exception:
        return
    rows = d.get("rows", d.get("guidelines", []))
    for ri, r in enumerate(rows[:tbl.rows.__len__() - 1]):
        data = [_str(r.get("guideline",""))[:60], _str(r.get("line",""))[:40],
                _str(r.get("recommendation",""))[:100], _str(r.get("source",""))[:50]]
        for ci in range(min(len(data), tbl.columns.__len__())):
            try:
                cell = tbl.cell(ri + 1, ci)
                cell.text = data[ci]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)
            except Exception:
                pass
    for ri in range(len(rows) + 1, tbl.rows.__len__()):
        for ci in range(tbl.columns.__len__()):
            try:
                tbl.cell(ri, ci).text = ""
            except Exception:
                pass


def _fill_treatment_algo(slide, d: dict, meta: dict):
    """Slide 12: Fill 13x3 table_area. Col 1 stays empty (Ovals overlap it)."""
    _set_text(_find_shape(slide, "header_title"),
              f"Treatment Landscape — {meta.get('country','Global')}")
    table_shape = _find_shape(slide, "table_area")
    if table_shape is None:
        return
    try:
        tbl = table_shape.table
    except Exception:
        return
    lines = d.get("lines", [])
    flat_rows = []
    for line in lines:
        regs = line.get("regimens", [])
        if isinstance(regs, str):
            regs = [regs]
        for reg in regs[:8]:
            if isinstance(reg, dict):
                name = _str(reg.get("name", reg.get("regimen", "")))
                note = _str(reg.get("note", reg.get("approval_status", "")))
                flat_rows.append([_str(line.get("line", "")), "", f"       {name} — {note}"])
            else:
                flat_rows.append([_str(line.get("line", "")), "", f"       {_str(reg)}"])
    for ri, row in enumerate(flat_rows[:tbl.rows.__len__() - 1]):
        for ci in range(min(len(row), tbl.columns.__len__())):
            try:
                cell = tbl.cell(ri + 1, ci)
                cell.text = row[ci][:100]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9)
            except Exception:
                pass
    for ri in range(len(flat_rows) + 1, tbl.rows.__len__()):
        for ci in range(tbl.columns.__len__()):
            try:
                tbl.cell(ri, ci).text = ""
            except Exception:
                pass


def _fill_imperatives(slide, d: dict, meta: dict):
    """Slide 11: Rectangle 39/80/81/82 (pillars) + Rectangle 53 (base bar)."""
    _set_text(_find_shape(slide, "header_title"),
              f"Strategic Imperatives — {meta.get('country','Global')} {meta.get('year','')}")
    pillars = d.get("pillars", [])
    pillar_shapes = ["Rectangle 39", "Rectangle 80", "Rectangle 81", "Rectangle 82"]
    for i, pillar in enumerate(pillars[:4]):
        if i < len(pillar_shapes):
            shape = _find_shape(slide, pillar_shapes[i])
            if shape:
                title = _str(pillar.get("title", ""))
                objs = "\n".join(f"• {_str(o)[:80]}" for o in pillar.get("objectives", [])[:4])
                _set_text(shape, f"{title}\n\n{objs}")
    base = _find_shape(slide, "Rectangle 53")
    if base:
        _set_text(base, f"MEDICAL AFFAIRS PLAN — {meta.get('drug','')} — {meta.get('year','')}")


def _fill_tactics(slide, d: dict, meta: dict):
    """Slide 13: 3 categories (Rectangle 2/8/9), initiatives (7/10/11), outcomes (12/13/14)."""
    _set_text(_find_shape(slide, "header_title"), f"Tactical Plan — {meta.get('year','')}")
    rows = d.get("rows", [])
    categories = {}
    for r in rows:
        cat = _str(r.get("type", "Other"))
        if cat not in categories:
            categories[cat] = []
        categories[cat].append(r)
    cat_list = list(categories.items())[:3]
    cat_titles = ["Rectangle 2", "Rectangle 8", "Rectangle 9"]
    cat_inits = ["Rectangle 7", "Rectangle 10", "Rectangle 11"]
    cat_outcomes = ["Rectangle 12", "Rectangle 13", "Rectangle 14"]
    for i, (cat_name, cat_rows) in enumerate(cat_list):
        if i < len(cat_titles):
            _set_text(_find_shape(slide, cat_titles[i]), cat_name)
        if i < len(cat_inits):
            inits = "\n".join(f"• {_str(r.get('tactic',''))[:50]}" for r in cat_rows[:4])
            _set_text(_find_shape(slide, cat_inits[i]), inits)
        if i < len(cat_outcomes):
            kpis = "\n".join(_str(r.get("kpi", r.get("description","")))[:40] for r in cat_rows[:2])
            _set_text(_find_shape(slide, cat_outcomes[i]), kpis or "KPIs TBD")


def _fill_narrative(slide, d: dict, meta: dict):
    """Slide 14: Rectangle 19 (summary), findings_body, reco_body."""
    _set_text(_find_shape(slide, "header_title"), f"Scientific Narrative — {meta.get('drug','')}")
    _set_text(_find_shape(slide, "Rectangle 19"),
              _str(d.get("primary_message", ""))[:200])
    tps = d.get("talking_points", [])
    evidence = "\n\n".join(
        f"• {_str(tp.get('focus',''))}: {', '.join(_str(p)[:50] for p in tp.get('points',[])[:3])}"
        for tp in tps[:4]
    )
    _set_text(_find_shape(slide, "findings_body"), evidence or _str(d.get("primary_message", "")))
    positioning = d.get("competitive_context", d.get("key_evidence_statement", ""))
    se = d.get("supporting_evidence", "")
    reco = f"{positioning}\n\n{se}" if positioning and se else positioning or se
    _set_text(_find_shape(slide, "reco_body"), reco or "See evidence section.")


def _fill_unmet_needs(slide, d: dict, meta: dict):
    """Slide 15: 3 challenge/solution pairs."""
    _set_text(_find_shape(slide, "header_title"),
              f"Unmet Medical Needs — {meta.get('country','Global')}")
    needs = d.get("needs", [])
    challenge_shapes = ["TextBox 20", "TextBox 3", "TextBox 19"]
    solution_shapes = ["TextBox 11", "TextBox 4", "TextBox 24"]
    for i in range(3):
        if i < len(needs):
            n = needs[i]
            mag = n.get("magnitude", "")
            icon = "🔴" if mag == "HIGH" else "🟡" if mag == "MEDIUM" else "🟢"
            _set_text(_find_shape(slide, challenge_shapes[i]) if i < len(challenge_shapes) else None,
                      f"{icon} {_str(n.get('title',''))}")
            _set_text(_find_shape(slide, solution_shapes[i]) if i < len(solution_shapes) else None,
                      _str(n.get("detail", ""))[:120])
        else:
            if i < len(challenge_shapes):
                _set_text(_find_shape(slide, challenge_shapes[i]), "")
            if i < len(solution_shapes):
                _set_text(_find_shape(slide, solution_shapes[i]), "")


def _fill_timeline(slide, d: dict, meta: dict):
    """Slide 16: 6x 'Rounded Rectangle 27' milestones."""
    _set_text(_find_shape(slide, "header_title"),
              f"Medical Affairs Roadmap — {meta.get('country','Global')} {meta.get('year','')}")
    events = d.get("events", [])
    milestones = _find_shapes_by_name(slide, "Rounded Rectangle 27")
    for i, ms in enumerate(milestones[:6]):
        if i < len(events):
            ev = events[i]
            _set_text(ms, f"{_str(ev.get('quarter',''))}: {_str(ev.get('event',''))}\n{_str(ev.get('detail',''))}"[:80])
        else:
            _set_text(ms, "")


def _fill_market_access(slide, d: dict, meta: dict):
    """Slide 17: Center (TextBox 8/9) + 5 country TextBoxes around it."""
    _set_text(_find_shape(slide, "header_title"), f"Market Access — {meta.get('country','Global')}")
    _set_text(_find_shape(slide, "Subtitle 2"), "Regulatory & Reimbursement Status")
    ap = d.get("approval_status", {})
    _set_text(_find_shape(slide, "TextBox 8"), meta.get("drug", ""))
    center = _str(ap.get("national_authority", ap.get("ema", "")))[:80]
    _set_text(_find_shape(slide, "TextBox 9"), center)
    country_shapes = ["TextBox 34", "TextBox 22", "TextBox 19", "TextBox 20", "TextBox 21"]
    country_data = []
    if ap.get("ema"):
        country_data.append(f"EMA\n{_str(ap['ema'])[:50]}")
    if ap.get("fda"):
        country_data.append(f"FDA\n{_str(ap['fda'])[:50]}")
    if ap.get("national_authority"):
        country_data.append(f"{meta.get('country','')}\n{_str(ap['national_authority'])[:50]}")
    for r in d.get("reimbursement_table", [])[:3]:
        country_data.append(f"{_str(r.get('body', r.get('country','')))[:20]}\n{_str(r.get('decision', r.get('status','')))[:50]}")
    for i, cs in enumerate(country_shapes):
        shape = _find_shape(slide, cs)
        if shape and i < len(country_data):
            _set_text(shape, country_data[i])
        elif shape:
            _set_text(shape, "")


def _fill_swot(slide, d: dict, meta: dict):
    """Slide 18: Titles stay, fill 4x 'Oval 40' with bullets."""
    _set_text(_find_shape(slide, "header_title"),
              f"SWOT Analysis — {meta.get('drug','')} {meta.get('country','Global')}")
    ovals = _find_shapes_by_name(slide, "Oval 40")
    for i, key in enumerate(["strengths", "weaknesses", "opportunities", "threats"]):
        items = d.get(key, [])
        bullets = "\n".join(f"• {_str(item)[:55]}" for item in items[:5])
        if i < len(ovals):
            _set_text(ovals[i], bullets or f"No {key}")


def _fill_competitor_table(slide, d: dict, meta: dict):
    """Slide 19: Fill header + source. Chart/table rendered as extra slides."""
    _set_text(_find_shape(slide, "header_title"),
              f"Competitive Landscape — {meta.get('country','Global')} {meta.get('year','')}")
    _set_text(_find_shape(slide, "source_text"),
              "Source: EMA EPAR, FDA, NCCN, Published Phase 3 data")
    rows = d.get("rows", [])
    focus = [r for r in rows if r.get("is_focus")]
    if focus:
        hr = _str(focus[0].get("hr_pfs", ""))
        if hr:
            _set_text(_find_shape(slide, "HR Annotation"), f"HR: {hr}")


# ══════════════════════════════════════════════════════════
# Dispatch: section_type → template fill function
# Covers ALL 20 section types (was 10, now 20)
# ══════════════════════════════════════════════════════════
TEMPLATE_FILLERS = {
    "executive_summary": _fill_executive_summary,
    "disease_intro": _fill_disease_intro,
    "prevalence_kpi": _fill_prevalence_kpi,
    "moa": _fill_moa,
    "pivotal_table": _fill_pivotal_table,
    "differentiators": _fill_differentiators,
    "areas_interest": _fill_areas_interest,
    "iep": _fill_iep,
    "guidelines": _fill_guidelines,
    "treatment_algo": _fill_treatment_algo,
    "swot": _fill_swot,
    "strategic_imperatives": _fill_imperatives,
    "imperatives": _fill_imperatives,
    "narrative": _fill_narrative,
    "differentiators": _fill_differentiators,
    "competitor_table": _fill_competitor_table,
    "unmet_needs": _fill_unmet_needs,
    "tactics": _fill_tactics,
    "tactical_plan": _fill_tactics,
    "timeline": _fill_timeline,
    "market_access": _fill_market_access,
}


# ══════════════════════════════════════════════════════════
# THEME PALETTES (hex without #)
# ══════════════════════════════════════════════════════════
PALETTES = {
    "dark":    {"bg":"0B1A3B","navy":"0D2B4E","surface":"163060","white":"FFFFFF","text":"EAF0FF","muted":"7B9FD4","dim":"4A6A9A","accent":"7C6FFF","teal":"22D3A5","gold":"F5C842","rose":"FF5F7E"},
    "light":   {"bg":"F0F4F8","navy":"FFFFFF","surface":"E2E8F0","white":"1A202C","text":"2D3748","muted":"718096","dim":"A0AEC0","accent":"3182CE","teal":"319795","gold":"D69E2E","rose":"E53E3E"},
    "gray":    {"bg":"2C3E50","navy":"34495E","surface":"415B73","white":"FFFFFF","text":"ECF0F1","muted":"BDC3C7","dim":"7F8C8D","accent":"E74C3C","teal":"1ABC9C","gold":"F39C12","rose":"E74C3C"},
    "pharma":  {"bg":"E8F4FD","navy":"F0F9FF","surface":"DBEAFE","white":"1E3A5F","text":"1E3A5F","muted":"4A6FA5","dim":"93B4D4","accent":"2563EB","teal":"0D9488","gold":"D97706","rose":"DC2626"},
    "premium": {"bg":"0A0A0A","navy":"111111","surface":"1C1C1C","white":"FFFFFF","text":"E8E8E8","muted":"999999","dim":"555555","accent":"C9A84C","teal":"C9A84C","gold":"C9A84C","rose":"B33030"},
    # Actual deployed themes
    "normal": {"bg":"0B1A3B","navy":"0D2B4E","surface":"163060","white":"FFFFFF","text":"EAF0FF","muted":"7B9FD4","dim":"4A6A9A","accent":"7C6FFF","teal":"22D3A5","gold":"F5C842","rose":"FF5F7E"},
    "gold": {"bg":"0A0A0A","navy":"111111","surface":"1C1C1C","white":"FFFFFF","text":"E8E8E8","muted":"999999","dim":"555555","accent":"C9A84C","teal":"C9A84C","gold":"D4AF37","rose":"B33030"},
    "aquarell": {"bg":"F0F4F8","navy":"FFFFFF","surface":"E2E8F0","white":"1A202C","text":"2D3748","muted":"718096","dim":"A0AEC0","accent":"3182CE","teal":"319795","gold":"D69E2E","rose":"E53E3E"},
}

FONT = "Calibri"

# Slideworks 2025 theme color resolution map
# Used by _get_text_color_deep to resolve schemeClr references
SLIDEWORKS_THEME = {
    "dk1": "000000", "lt1": "FFFFFF", "dk2": "000000", "lt2": "FFFFFF",
    "tx1": "000000", "tx2": "000000", "bg1": "FFFFFF", "bg2": "FFFFFF",
    "accent1": "051C2C", "accent2": "22A3DF", "accent3": "2222F6",
    "accent4": "1877A2", "accent5": "A2D5EC", "accent6": "F8BC3B",
}
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MX = Inches(0.5)       # left margin
MW = Inches(12.33)      # content width
HEADER_H = Inches(0.82)


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _str(v) -> str:
    """Safe string extraction (mirrors frontend str() helper)."""
    if isinstance(v, str):
        return v
    if isinstance(v, dict):
        return v.get("bullet") or v.get("text") or v.get("title") or str(v)
    return str(v) if v else ""


def _add_textbox(slide, left, top, width, height, text,
                 font_size=10, font_color="EAF0FF", bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font_name=FONT,
                 word_wrap=True):
    """Add a textbox with consistent formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(font_color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    tf.auto_size = None
    return txBox


def _add_rounded_rect(slide, left, top, width, height, fill_color, radius=Inches(0.1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()
    # Adjust corner radius via XML
    if hasattr(shape, '_element'):
        sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if sp_pr is not None:
            av_lst = sp_pr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
            if av_lst is not None:
                for gd in av_lst:
                    gd.set('fmla', f'val {int(radius / Inches(1) * 10000)}')
    return shape


def _add_rect(slide, left, top, width, height, fill_color):
    """Add a simple rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════
# SLIDE HELPERS (header, footer, divider)
# ══════════════════════════════════════════════════════════
def _slide_header(slide, title: str, P: dict, subtitle: str = ""):
    """Add themed header bar to a slide."""
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, P["navy"])
    _add_rect(slide, Inches(0), Inches(0), Inches(0.15), HEADER_H, P["teal"])
    _add_textbox(slide, Inches(0.28), Inches(0), Inches(9.5), HEADER_H,
                 title, font_size=15, font_color=P["white"], bold=True,
                 valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_textbox(slide, Inches(10), Inches(0), Inches(3.1), HEADER_H,
                     subtitle, font_size=9, font_color=P["teal"],
                     alignment=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def _slide_footer(slide, P: dict, refs: list[str] = None):
    """Add footer with disclaimer and references."""
    _add_textbox(slide, MX, Inches(7.08), MW, Inches(0.30),
                 "⚠ AI-generated — Verify before external use · MedAI Suite Premium",
                 font_size=7, font_color=P["dim"], italic=True)
    if refs:
        clean = [r for r in refs if r][:6]
        if clean:
            txt = "  ·  ".join(f"[{i+1}] {r}" for i, r in enumerate(clean))
            _add_textbox(slide, MX, Inches(6.82), MW, Inches(0.24),
                         txt, font_size=7, font_color=P["dim"], italic=True)


def _add_divider_slide(pres, num: int, title: str, P: dict):
    """Add a section divider slide."""
    sl = pres.slides.add_slide(pres.slide_layouts[6])  # Blank
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = _rgb(P["bg"])
    _add_rect(sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, P["bg"])
    _add_textbox(sl, MX, Inches(1.5), Inches(3), Inches(2.5),
                 str(num).zfill(2), font_size=72, font_color=P["accent"], bold=True,
                 valign=MSO_ANCHOR.MIDDLE)
    _add_textbox(sl, Inches(3.8), Inches(2.0), Inches(8), Inches(1.5),
                 title, font_size=28, font_color=P["white"], bold=True,
                 valign=MSO_ANCHOR.MIDDLE)
    _add_rect(sl, Inches(3.8), Inches(3.6), Inches(2), Inches(0.04), P["teal"])


def _new_slide(pres, P: dict):
    """Create a new blank slide with background fill."""
    sl = pres.slides.add_slide(pres.slide_layouts[6])  # Blank layout
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = _rgb(P["navy"])
    return sl


def _card_grid(slide, cards: list[dict], y: float, h: float, P: dict):
    """Add a grid of info cards (up to 6, in rows of 3)."""
    y_in = Inches(y)
    h_in = Inches(h)
    n = min(len(cards), 6)
    cols = min(n, 3)
    cw = (12.33 - 0.3 * (cols - 1)) / cols

    for i, card in enumerate(cards[:6]):
        col = i % cols
        row = i // cols
        cx = 0.5 + col * (cw + 0.3)
        cy = y + row * (h + 0.15)

        _add_rounded_rect(slide, Inches(cx), Inches(cy), Inches(cw), Inches(h), P["surface"])
        _add_textbox(slide, Inches(cx + 0.12), Inches(cy + 0.08), Inches(cw - 0.24), Inches(0.35),
                     card.get("title", ""), font_size=10, font_color=P["teal"], bold=True)
        _add_textbox(slide, Inches(cx + 0.12), Inches(cy + 0.42), Inches(cw - 0.24), Inches(h - 0.55),
                     card.get("body", ""), font_size=9, font_color=P["text"])


# ══════════════════════════════════════════════════════════
# EXTRACT REFERENCES from section data
# ══════════════════════════════════════════════════════════
def _extract_refs(d: dict) -> list[str]:
    refs = set()
    def walk(obj):
        if isinstance(obj, dict):
            for k, v in obj.items():
                if k in ("source", "key_reference") and isinstance(v, str) and len(v) > 5:
                    refs.add(v)
                elif k == "references" and isinstance(v, list):
                    for r in v:
                        if isinstance(r, str):
                            refs.add(r)
                        elif isinstance(r, dict) and "text" in r:
                            refs.add(r["text"])
                else:
                    walk(v)
        elif isinstance(obj, list):
            for item in obj:
                walk(item)
    walk(d)
    return list(refs)


# ══════════════════════════════════════════════════════════
# SECTION RENDERERS — one function per section_type
# ══════════════════════════════════════════════════════════

def render_executive_summary(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Executive Summary — {meta.get('country','Global')} {meta.get('year','')}", P)
    rows = d.get("rows", [])
    _card_grid(sl, [{"title": f"{r.get('icon','')} {r.get('topic','')}", "body": r.get("summary", "")} for r in rows], 1.1, 1.8, P)
    _slide_footer(sl, P, _extract_refs(d))


def render_prevalence_kpi(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Prevalence & Epidemiology — {meta.get('country','Global')}", P)
    kpis = d.get("kpis", [])
    for i, k in enumerate(kpis[:4]):
        x = 0.5 + i * 3.1
        _add_rounded_rect(sl, Inches(x), Inches(1.2), Inches(2.9), Inches(1.4), P["surface"])
        _add_textbox(sl, Inches(x), Inches(1.25), Inches(2.9), Inches(0.7),
                     _str(k.get("value", "")), font_size=22, font_color=P["teal"], bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(sl, Inches(x), Inches(1.95), Inches(2.9), Inches(0.5),
                     _str(k.get("label", "")), font_size=8, font_color=P["muted"],
                     alignment=PP_ALIGN.CENTER)
    # Context table
    ct = d.get("context_table", [])
    if ct:
        from pptx.util import Inches as In
        tbl_data = [["Metric", "Value", "Source"]] + [[_str(r.get("metric","")), _str(r.get("value","")), _str(r.get("source",""))] for r in ct[:8]]
        rows_n = len(tbl_data)
        cols_n = 3
        tbl_shape = sl.shapes.add_table(rows_n, cols_n, MX, Inches(2.9), MW, Inches(min(3.5, rows_n * 0.4)))
        tbl = tbl_shape.table
        col_widths = [Inches(5), Inches(3), Inches(4.33)]
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w
        for ri, row in enumerate(tbl_data):
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = cell_text
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(9 if ri > 0 else 8)
                        run.font.color.rgb = _rgb(P["text"] if ri > 0 else P["muted"])
                        run.font.name = FONT
                        if ri == 0:
                            run.font.bold = True
    _slide_footer(sl, P, _extract_refs(d))


def render_treatment_algo(pres, d: dict, P: dict, meta: dict):
    lines = d.get("lines", [])
    for line in lines:
        sl = _new_slide(pres, P)
        line_name = _str(line.get("line", ""))
        _slide_header(sl, f"Treatment Landscape — {line_name}", P, meta.get("country", "Global"))
        regs = line.get("regimens", [])
        if isinstance(regs, str):
            regs = [regs]
        # Build table
        tbl_data = [["Regimen", "Details"]]
        for r in regs[:10]:
            txt = _str(r)
            tbl_data.append([txt[:60], txt[60:120] if len(txt) > 60 else ""])
        if len(tbl_data) > 1:
            rows_n = len(tbl_data)
            tbl_shape = sl.shapes.add_table(rows_n, 2, MX, Inches(1.1), MW, Inches(min(5.5, rows_n * 0.5)))
            tbl = tbl_shape.table
            tbl.columns[0].width = Inches(6)
            tbl.columns[1].width = Inches(6.33)
            for ri, row in enumerate(tbl_data):
                for ci, cell_text in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = cell_text
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(9)
                            run.font.color.rgb = _rgb(P["text"] if ri > 0 else P["muted"])
                            run.font.name = FONT
        _slide_footer(sl, P, _extract_refs(d))


def render_pivotal_table(pres, d: dict, P: dict, meta: dict):
    trials = d.get("trials", [])
    active = [t for t in trials if t.get("include_in_map") is not False]

    for tr in active:
        sl = _new_slide(pres, P)
        eff = tr.get("efficacy", {})
        _slide_header(sl, f"{_str(tr.get('name',''))} — {_str(tr.get('phase',''))}", P, _str(tr.get("design", "")))

        # KPI boxes
        metrics = [
            ("mPFS", _str(eff.get("mpfs_drug", "")), _str(eff.get("mpfs_control", ""))),
            ("HR", _str(eff.get("hr_pfs", "")), ""),
            ("ORR", _str(eff.get("orr_drug", "")), _str(eff.get("orr_control", ""))),
            ("mOS", _str(eff.get("mos_drug", "")), _str(eff.get("mos_control", ""))),
        ]
        for i, (lbl, val, vs) in enumerate(metrics):
            x = 0.5 + i * 3.1
            _add_rounded_rect(sl, Inches(x), Inches(1.2), Inches(2.9), Inches(1.5), P["surface"])
            _add_textbox(sl, Inches(x), Inches(1.25), Inches(2.9), Inches(0.3),
                         lbl, font_size=9, font_color=P["muted"], alignment=PP_ALIGN.CENTER,
                         font_name="DM Mono")
            _add_textbox(sl, Inches(x), Inches(1.55), Inches(2.9), Inches(0.6),
                         val or "N/A", font_size=20, font_color=P["teal"], bold=True,
                         alignment=PP_ALIGN.CENTER)
            if vs:
                _add_textbox(sl, Inches(x), Inches(2.2), Inches(2.9), Inches(0.3),
                             f"vs {vs}", font_size=9, font_color=P["dim"], alignment=PP_ALIGN.CENTER)

        # KM Curve — choose best endpoint: prefer specific data per trial
        theme_key = _theme_key_from_palette(P)
        km_rendered = False
        
        # Try PFS first
        pfs_drug = _parse_months(eff.get("mpfs_drug", ""))
        pfs_ctrl = _parse_months(eff.get("mpfs_control", ""))
        if pfs_drug and pfs_ctrl and pfs_drug > 0 and pfs_ctrl > 0:
            km_png = chart_renderer.render_km_curve(
                drug_name=meta.get("drug", "Drug"),
                drug_median=pfs_drug,
                control_median=pfs_ctrl,
                endpoint="PFS",
                n_total=_str(tr.get("n_total", "N/A")),
                hr=_str(eff.get("hr_pfs", "")),
                p_value=_str(eff.get("p_value_pfs", eff.get("p_value", ""))),
                theme=theme_key,
            )
            if km_png:
                img_stream = io.BytesIO(km_png)
                sl.shapes.add_picture(img_stream, Inches(6.5), Inches(3.2), Inches(6.0), Inches(3.5))
                km_rendered = True
        
        # Also try OS — if different from PFS, add as second chart below safety
        os_drug = _parse_months(eff.get("mos_drug", ""))
        os_ctrl = _parse_months(eff.get("mos_control", ""))
        if os_drug and os_ctrl and os_drug > 0 and os_ctrl > 0:
            # Only add if OS data is actually different from PFS
            if not km_rendered or abs(os_drug - (pfs_drug or 0)) > 1.0:
                os_png = chart_renderer.render_km_curve(
                    drug_name=meta.get("drug", "Drug"),
                    drug_median=os_drug,
                    control_median=os_ctrl,
                    endpoint="OS",
                    n_total=_str(tr.get("n_total", "N/A")),
                    hr=_str(eff.get("hr_os", "")),
                    p_value=_str(eff.get("p_value_os", "")),
                    theme=theme_key,
                )
                if os_png:
                    img_stream = io.BytesIO(os_png)
                    if km_rendered:
                        # Put OS on a separate slide
                        sl2 = _new_slide(pres, P)
                        _slide_header(sl2, f"{_str(tr.get('name',''))} — Overall Survival", P, _str(tr.get("design", "")))
                        sl2.shapes.add_picture(img_stream, Inches(1.5), Inches(1.5), Inches(10.0), Inches(5.0))
                        _slide_footer(sl2, P, [tr.get("source", "")])
                    else:
                        sl.shapes.add_picture(img_stream, Inches(6.5), Inches(3.2), Inches(6.0), Inches(3.5))

        # Safety text
        saf = tr.get("safety", {})
        saf_lines = []
        if saf.get("grade34_heme"):
            saf_lines.append("Grade 3-4 Heme: " + ", ".join(_str(s) for s in saf["grade34_heme"]))
        if saf.get("grade34_nonheme"):
            saf_lines.append("Grade 3-4 Non-heme: " + ", ".join(_str(s) for s in saf["grade34_nonheme"]))
        if saf.get("discontinuation_rate"):
            saf_lines.append(f"Discontinuation: {_str(saf['discontinuation_rate'])}")
        if saf_lines:
            _add_textbox(sl, MX, Inches(3.5), Inches(6.0), Inches(1.5),
                         "\n".join(saf_lines), font_size=9, font_color=P["text"])

        _slide_footer(sl, P, [tr.get("source", "")])


def render_competitor_table(pres, d: dict, P: dict, meta: dict):
    rows = d.get("rows", [])
    theme_key = _theme_key_from_palette(P)
    
    # mPFS Bar chart
    if len(rows) >= 3:
        chart_png = chart_renderer.render_mpfs_bar_chart(
            rows=rows, focus_drug=meta.get("drug", ""), theme=theme_key,
        )
        if chart_png:
            sl = _new_slide(pres, P)
            _slide_header(sl, "Competitive Landscape — mPFS Comparison", P, f"{meta.get('country','Global')} {meta.get('year','')}")
            img_stream = io.BytesIO(chart_png)
            sl.shapes.add_picture(img_stream, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))
            _slide_footer(sl, P, _extract_refs(d))

    # ORR Bar chart (NEW)
    if len(rows) >= 3:
        orr_png = chart_renderer.render_orr_bar_chart(
            rows=rows, focus_drug=meta.get("drug", ""), theme=theme_key,
        )
        if orr_png:
            sl = _new_slide(pres, P)
            _slide_header(sl, "Competitive Landscape — ORR Comparison", P, f"{meta.get('country','Global')} {meta.get('year','')}")
            img_stream = io.BytesIO(orr_png)
            sl.shapes.add_picture(img_stream, Inches(0.3), Inches(1.0), Inches(12.7), Inches(5.8))
            _slide_footer(sl, P, _extract_refs(d))

    # Table slides
    per_page = 7
    for pg in range(0, len(rows), per_page):
        sl = _new_slide(pres, P)
        page_rows = rows[pg:pg + per_page]
        pgn = f" ({pg // per_page + 1}/{(len(rows) - 1) // per_page + 1})" if len(rows) > per_page else ""
        _slide_header(sl, f"Competitive Landscape{pgn}", P, meta.get("country", "Global"))

        tbl_data = [["Drug / Trial", "LOT", "mPFS", "HR", "ORR", "Key Differentiator"]]
        for r in page_rows:
            tbl_data.append([
                _str(r.get("drug_trial", "")),
                _str(r.get("prior_lot", "")),
                _str(r.get("mpfs", "")),
                _str(r.get("hr_pfs", "")),
                _str(r.get("orr", "")),
                _str(r.get("key_differentiator", r.get("notes", ""))),
            ])
        rows_n = len(tbl_data)
        tbl_shape = sl.shapes.add_table(rows_n, 6, MX, Inches(1.1), MW, Inches(min(5.5, rows_n * 0.55)))
        tbl = tbl_shape.table
        widths = [Inches(3), Inches(0.8), Inches(1.8), Inches(1.8), Inches(1.0), Inches(3.93)]
        for ci, w in enumerate(widths):
            tbl.columns[ci].width = w
        for ri, row in enumerate(tbl_data):
            for ci, text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = text[:100]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(8)
                        run.font.color.rgb = _rgb(P["text"] if ri > 0 else P["muted"])
                        run.font.name = FONT
        _slide_footer(sl, P, _extract_refs(d))


def render_swot(pres, d: dict, P: dict, meta: dict):
    # Try chart version
    theme_key = _theme_key_from_palette(P)
    chart_png = chart_renderer.render_swot_chart(
        strengths=[_str(s) for s in d.get("strengths", [])],
        weaknesses=[_str(s) for s in d.get("weaknesses", [])],
        opportunities=[_str(s) for s in d.get("opportunities", [])],
        threats=[_str(s) for s in d.get("threats", [])],
        theme=theme_key,
    )
    sl = _new_slide(pres, P)
    _slide_header(sl, f"SWOT Analysis — {meta.get('drug','')} {meta.get('country','Global')} {meta.get('year','')}", P)
    if chart_png:
        img_stream = io.BytesIO(chart_png)
        sl.shapes.add_picture(img_stream, Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.0))
    _slide_footer(sl, P, _extract_refs(d))


def render_imperatives(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Strategic Imperatives — {meta.get('country','Global')} {meta.get('year','')}", P)
    pillars = d.get("pillars", [])
    n = min(len(pillars), 4)
    pillar_colors = [P["teal"], P["accent"], P["gold"], P["rose"]]

    # Base bar
    _add_rounded_rect(sl, Inches(0.3), Inches(6.25), Inches(12.73), Inches(0.35), P["surface"])
    _add_textbox(sl, Inches(0.3), Inches(6.25), Inches(12.73), Inches(0.35),
                 f"MEDICAL AFFAIRS PLAN — {meta.get('drug','')} — {meta.get('year','')}",
                 font_size=9, font_color=P["muted"], bold=True, alignment=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)

    # Roof
    _add_rounded_rect(sl, Inches(0.3), Inches(1.05), Inches(12.73), Inches(0.55), P["surface"])
    _add_rect(sl, Inches(0.3), Inches(1.56), Inches(12.73), Inches(0.04), P["teal"])
    _add_textbox(sl, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.55),
                 f"Strategic Vision: {meta.get('drug','')} — {meta.get('country','Global')}",
                 font_size=12, font_color=P["white"], bold=True, alignment=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)

    # Pillars
    if n > 0:
        gap = 0.25
        total_w = 12.33
        pw = (total_w - gap * (n - 1)) / n
        pillar_top = 1.68
        pillar_bot = 6.17
        p_h = pillar_bot - pillar_top

        for i, pillar in enumerate(pillars[:n]):
            px = 0.5 + i * (pw + gap)
            pc = pillar_colors[i % len(pillar_colors)]
            _add_rounded_rect(sl, Inches(px), Inches(pillar_top), Inches(pw), Inches(p_h), P["surface"])
            _add_rect(sl, Inches(px), Inches(pillar_top), Inches(pw), Inches(0.08), pc)
            _add_rect(sl, Inches(px), Inches(pillar_bot - 0.06), Inches(pw), Inches(0.06), pc)
            _add_textbox(sl, Inches(px + 0.1), Inches(pillar_top + 0.15), Inches(pw - 0.2), Inches(0.45),
                         _str(pillar.get("title", "")), font_size=11, font_color=pc, bold=True,
                         alignment=PP_ALIGN.CENTER)
            _add_rect(sl, Inches(px + 0.2), Inches(pillar_top + 0.65), Inches(pw - 0.4), Inches(0.02), pc)
            objs = pillar.get("objectives", [])[:4]
            obj_text = "\n\n".join(f"• {_str(o)[:90]}" for o in objs)
            _add_textbox(sl, Inches(px + 0.12), Inches(pillar_top + 0.75), Inches(pw - 0.24), Inches(p_h - 1.0),
                         obj_text, font_size=8, font_color=P["text"])

    _slide_footer(sl, P, _extract_refs(d))


def render_narrative(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Scientific Narrative — {meta.get('drug','')}", P)
    y = 1.1
    if d.get("primary_message"):
        _add_rounded_rect(sl, MX, Inches(y), MW, Inches(1.0), P["surface"])
        _add_rect(sl, MX, Inches(y), Inches(0.12), Inches(1.0), P["teal"])
        _add_textbox(sl, Inches(0.75), Inches(y + 0.1), Inches(11.8), Inches(0.8),
                     d["primary_message"], font_size=11, font_color=P["text"], bold=True)
        y += 1.2

    tps = d.get("talking_points", [])
    for tp in tps[:4]:
        _add_textbox(sl, MX, Inches(y), MW, Inches(0.3),
                     _str(tp.get("focus", "")), font_size=10, font_color=P["accent"], bold=True)
        points = tp.get("points", [])
        pts_text = "\n".join(f"• {_str(p)[:120]}" for p in points[:3])
        _add_textbox(sl, Inches(0.7), Inches(y + 0.3), Inches(12.0), Inches(0.8),
                     pts_text, font_size=9, font_color=P["text"])
        y += 1.2

    _slide_footer(sl, P, _extract_refs(d))


def render_timeline(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Medical Affairs Roadmap — {meta.get('country','Global')} {meta.get('year','')}", P)
    events = d.get("events", [])
    quarters = ["Q1", "Q2", "Q3", "Q4"]
    q_colors = [P["accent"], P["teal"], P["gold"], P["rose"]]

    # Road
    road_y = 3.55
    road_h = 0.7
    _add_rounded_rect(sl, Inches(0.2), Inches(road_y), Inches(12.93), Inches(road_h), "3A3F47")
    # Lane markings
    for dx in [x * 0.6 + 0.5 for x in range(21)]:
        if dx < 12.8:
            _add_rounded_rect(sl, Inches(dx), Inches(road_y + road_h / 2 - 0.02), Inches(0.35), Inches(0.04), "F5C842")

    seg_w = (12.93 - 0.4) / 4
    for qi, q in enumerate(quarters):
        mx = 0.2 + 0.2 + qi * seg_w + seg_w / 2
        qc = q_colors[qi]
        # Quarter badge
        _add_rounded_rect(sl, Inches(mx - 0.28), Inches(road_y - 0.5), Inches(0.56), Inches(0.55), qc)
        _add_textbox(sl, Inches(mx - 0.28), Inches(road_y - 0.5), Inches(0.56), Inches(0.35),
                     q, font_size=12, font_color=P["bg"], bold=True, alignment=PP_ALIGN.CENTER,
                     valign=MSO_ANCHOR.MIDDLE)

        q_evts = [e for e in events if e.get("quarter") == q][:2]
        base_x = 0.2 + 0.2 + qi * seg_w
        for ei, ev in enumerate(q_evts):
            card_w = seg_w - 0.15
            card_h = 1.3
            card_x = base_x + 0.05
            above = (ei == 0)
            card_y = road_y - 0.55 - card_h if above else road_y + road_h + 0.25
            _add_rounded_rect(sl, Inches(card_x), Inches(card_y), Inches(card_w), Inches(card_h), P["surface"])
            # Color bar
            bar_y = card_y + card_h - 0.04 if above else card_y
            _add_rect(sl, Inches(card_x), Inches(bar_y), Inches(card_w), Inches(0.04), qc)
            _add_textbox(sl, Inches(card_x + 0.1), Inches(card_y + 0.06), Inches(card_w - 0.2), Inches(0.35),
                         _str(ev.get("event", ""))[:30], font_size=8, font_color=qc, bold=True)
            _add_textbox(sl, Inches(card_x + 0.1), Inches(card_y + 0.38), Inches(card_w - 0.2), Inches(card_h - 0.55),
                         _str(ev.get("detail", ""))[:80], font_size=7, font_color=P["text"])

    _slide_footer(sl, P, _extract_refs(d))


def render_generic(pres, d: dict, P: dict, meta: dict, label: str):
    """Fallback renderer for any section type without a dedicated function."""
    sl = _new_slide(pres, P)
    _slide_header(sl, label, P)
    slides_data = d.get("slides", [{}])
    sl_data = slides_data[0] if slides_data else {}
    items = sl_data.get("items", [])
    if items:
        _card_grid(sl, [{"title": _str(it.get("label", "")), "body": _str(it.get("text", ""))} for it in items], 1.1, 2.2, P)
    else:
        # Show key-value pairs
        lines = []
        for k, v in d.items():
            if k == "section_type" or isinstance(v, (dict, list)):
                continue
            lines.append(f"{k}: {_str(v)}")
        _add_textbox(sl, MX, Inches(1.2), MW, Inches(5.5),
                     "\n\n".join(lines) or "No data for this section.",
                     font_size=10, font_color=P["text"])
    _slide_footer(sl, P, _extract_refs(d))


def render_guidelines(pres, d: dict, P: dict, meta: dict):
    rows = d.get("rows", d.get("guidelines", []))
    per_page = 6
    for pg in range(0, max(len(rows), 1), per_page):
        sl = _new_slide(pres, P)
        pgn = f" ({pg // per_page + 1}/{(len(rows) - 1) // per_page + 1})" if len(rows) > per_page else ""
        _slide_header(sl, f"Guidelines — {meta.get('drug','')}{pgn}", P)
        page_rows = rows[pg:pg + per_page]
        if page_rows:
            tbl_data = [["Guideline", "Line", "Recommendation", "Source"]]
            for r in page_rows:
                tbl_data.append([_str(r.get("guideline",""))[:50], _str(r.get("line","")), _str(r.get("recommendation",""))[:120], _str(r.get("source",""))[:50]])
            rows_n = len(tbl_data)
            tbl_shape = sl.shapes.add_table(rows_n, 4, MX, Inches(1.1), MW, Inches(min(5.5, rows_n * 0.7)))
            tbl = tbl_shape.table
            for ci, w in enumerate([Inches(3), Inches(1.8), Inches(4.5), Inches(3.03)]):
                tbl.columns[ci].width = w
            for ri, row in enumerate(tbl_data):
                for ci, text in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = text
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(8 if ri == 0 else 9)
                            run.font.color.rgb = _rgb(P["muted"] if ri == 0 else P["text"])
                            run.font.name = FONT
        _slide_footer(sl, P, _extract_refs(d))


def render_unmet_needs(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Unmet Medical Needs — {meta.get('country','Global')}", P)
    needs = d.get("needs", [])
    _card_grid(sl, [{"title": ("🔴 " if n.get("magnitude") == "HIGH" else "🟡 ") + _str(n.get("title", "")), "body": _str(n.get("detail", ""))} for n in needs], 1.1, 2.2, P)
    _slide_footer(sl, P, _extract_refs(d))


def render_moa(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Mechanism of Action: {meta.get('drug','')}", P)
    if d.get("drug_class"):
        _add_textbox(sl, MX, Inches(1.1), MW, Inches(0.4),
                     f"{_str(d['drug_class'])} · Target: {_str(d.get('target',''))}",
                     font_size=11, font_color=P["teal"], bold=True)
    steps = d.get("pathway_steps", [])
    for i, s in enumerate(steps[:4]):
        y = 1.7 + i * 1.3
        _add_rounded_rect(sl, MX, Inches(y), Inches(0.5), Inches(0.5), P["accent"])
        _add_textbox(sl, MX, Inches(y), Inches(0.5), Inches(0.5),
                     str(s.get("step", i + 1)), font_size=14, font_color=P["white"], bold=True,
                     alignment=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_textbox(sl, Inches(1.2), Inches(y), Inches(3), Inches(0.4),
                     _str(s.get("title", "")), font_size=11, font_color=P["teal"], bold=True)
        _add_textbox(sl, Inches(1.2), Inches(y + 0.38), Inches(11), Inches(0.7),
                     _str(s.get("description", "")), font_size=9, font_color=P["text"])
    _slide_footer(sl, P, _extract_refs(d))


def render_differentiators(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Key Differentiators — {meta.get('drug','')}", P)
    diffs = d.get("differentiators", [])
    _card_grid(sl, [{"title": _str(df.get("title", "")), "body": _str(df.get("detail", df.get("evidence", "")))} for df in diffs], 1.1, 2.2, P)
    _slide_footer(sl, P, _extract_refs(d))


def render_market_access(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, f"Market Access — {meta.get('country','Global')}", P, "Regulatory & Reimbursement")
    ap = d.get("approval_status", {})
    agencies = [
        ("EMA", ap.get("ema", "")),
        ("FDA", ap.get("fda", "")),
        (meta.get("country", "National"), ap.get("national_authority", "")),
    ]
    for i, (lbl, val) in enumerate(a for a in agencies if a[1]):
        x = 0.5 + i * 4.15
        _add_rounded_rect(sl, Inches(x), Inches(1.15), Inches(3.95), Inches(1.8), P["surface"])
        _add_rect(sl, Inches(x), Inches(1.15), Inches(3.95), Inches(0.42), P["teal"])
        _add_textbox(sl, Inches(x + 0.1), Inches(1.15), Inches(2.5), Inches(0.42),
                     lbl, font_size=11, font_color=P["bg"], bold=True, valign=MSO_ANCHOR.MIDDLE)
        _add_textbox(sl, Inches(x + 0.15), Inches(1.65), Inches(3.65), Inches(1.2),
                     _str(val)[:200], font_size=8, font_color=P["text"])
    _slide_footer(sl, P, _extract_refs(d))


def render_tactics(pres, d: dict, P: dict, meta: dict):
    rows = d.get("rows", [])
    per_page = 5
    for pg in range(0, max(len(rows), 1), per_page):
        sl = _new_slide(pres, P)
        pgn = f" ({pg // per_page + 1}/{(len(rows) - 1) // per_page + 1})" if len(rows) > per_page else ""
        _slide_header(sl, f"Tactical Plan {meta.get('year','')}{pgn}", P, meta.get("country", "Global"))
        page_rows = rows[pg:pg + per_page]
        if page_rows:
            tbl_data = [["Type", "Tactic", "Description / KPI"]]
            for r in page_rows:
                tbl_data.append([_str(r.get("type","")), _str(r.get("tactic",""))[:80], _str(r.get("description", r.get("kpi","")))[:120]])
            rows_n = len(tbl_data)
            tbl_shape = sl.shapes.add_table(rows_n, 3, MX, Inches(1.1), MW, Inches(min(5.5, rows_n * 0.8)))
            tbl = tbl_shape.table
            for ci, w in enumerate([Inches(2), Inches(4), Inches(6.33)]):
                tbl.columns[ci].width = w
            for ri, row in enumerate(tbl_data):
                for ci, text in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = text
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(8 if ri == 0 else 9)
                            run.font.color.rgb = _rgb(P["muted"] if ri == 0 else P["text"])
                            run.font.name = FONT
        _slide_footer(sl, P, _extract_refs(d))


def render_subgroup_analysis(pres, d: dict, P: dict, meta: dict):
    subs = d.get("subgroups", [])
    per_page = 10
    for pg in range(0, max(len(subs), 1), per_page):
        sl = _new_slide(pres, P)
        pgn = f" ({pg // per_page + 1}/{(len(subs) - 1) // per_page + 1})" if len(subs) > per_page else ""
        _slide_header(sl, f"Subgroup Analysis — {meta.get('drug','')}{pgn}", P)
        page_subs = subs[pg:pg + per_page]
        if page_subs:
            tbl_data = [["Trial", "Subgroup", "Endpoint", "Drug", "HR", "Favours"]]
            for sg in page_subs:
                tbl_data.append([_str(sg.get("trial_name","")), _str(sg.get("subgroup","")), _str(sg.get("endpoint","")),
                                 _str(sg.get("mpfs_drug","")), _str(sg.get("hr","")), _str(sg.get("favours",""))])
            rows_n = len(tbl_data)
            tbl_shape = sl.shapes.add_table(rows_n, 6, MX, Inches(1.1), MW, Inches(min(5.5, rows_n * 0.45)))
            tbl = tbl_shape.table
            for ci, w in enumerate([Inches(2), Inches(2.8), Inches(1.3), Inches(1.8), Inches(2.5), Inches(1.93)]):
                tbl.columns[ci].width = w
            for ri, row in enumerate(tbl_data):
                for ci, text in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = text
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(8)
                            run.font.color.rgb = _rgb(P["text"] if ri > 0 else P["muted"])
                            run.font.name = FONT
        _slide_footer(sl, P, _extract_refs(d))


def render_areas_interest(pres, d: dict, P: dict, meta: dict):
    areas = d.get("areas", [])
    per_page = 2
    for pg in range(0, max(len(areas), 1), per_page):
        sl = _new_slide(pres, P)
        pgn = f" ({pg // per_page + 1}/{(len(areas) - 1) // per_page + 1})" if len(areas) > per_page else ""
        _slide_header(sl, f"Areas of Interest (ISR){pgn}", P, meta.get("drug", ""))
        page_areas = areas[pg:pg + per_page]
        for ai, area in enumerate(page_areas):
            y = 1.2 + ai * 2.8
            _add_rounded_rect(sl, MX, Inches(y), MW, Inches(2.5), P["surface"])
            _add_rect(sl, MX, Inches(y), Inches(0.1), Inches(2.5), P["teal"])
            _add_textbox(sl, Inches(0.75), Inches(y + 0.1), Inches(11.5), Inches(0.4),
                         _str(area.get("area", "")), font_size=12, font_color=P["teal"], bold=True)
            interests = area.get("interests", [])
            txt = "\n".join(f"• {_str(it)[:120]}" for it in interests[:5])
            _add_textbox(sl, Inches(0.75), Inches(y + 0.55), Inches(11.5), Inches(1.8),
                         txt, font_size=9, font_color=P["text"])
        _slide_footer(sl, P, _extract_refs(d))


def render_iep(pres, d: dict, P: dict, meta: dict):
    gaps = d.get("gaps", [])
    per_page = 6
    for pg in range(0, max(len(gaps), 1), per_page):
        sl = _new_slide(pres, P)
        pgn = f" ({pg // per_page + 1}/{(len(gaps) - 1) // per_page + 1})" if len(gaps) > per_page else ""
        _slide_header(sl, f"Integrated Evidence Plan{pgn}", P, f"{meta.get('drug','')} {meta.get('year','')}")
        page_gaps = gaps[pg:pg + per_page]
        if page_gaps:
            tbl_data = [["Evidence Gap", "Activity", "Responsible", "Status"]]
            for g in page_gaps:
                tbl_data.append([_str(g.get("gap",""))[:80], _str(g.get("activity",""))[:90], _str(g.get("responsible","")), _str(g.get("status",""))])
            rows_n = len(tbl_data)
            tbl_shape = sl.shapes.add_table(rows_n, 4, MX, Inches(1.1), MW, Inches(min(5.5, rows_n * 0.75)))
            tbl = tbl_shape.table
            for ci, w in enumerate([Inches(3.2), Inches(5.5), Inches(2), Inches(1.63)]):
                tbl.columns[ci].width = w
            for ri, row in enumerate(tbl_data):
                for ci, text in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = text
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(8)
                            run.font.color.rgb = _rgb(P["text"] if ri > 0 else P["muted"])
                            run.font.name = FONT
        _slide_footer(sl, P, _extract_refs(d))


def render_summary(pres, d: dict, P: dict, meta: dict):
    sl = _new_slide(pres, P)
    _slide_header(sl, "Summary & Next Steps", P)
    msgs = d.get("key_messages", [])
    for i, m in enumerate(msgs[:5]):
        y = 1.2 + i * 1.05
        _add_rounded_rect(sl, MX, Inches(y), MW, Inches(0.9), P["surface"])
        _add_rect(sl, MX, Inches(y), Inches(0.1), Inches(0.9), P["teal"])
        _add_textbox(sl, Inches(0.8), Inches(y), Inches(11.5), Inches(0.9),
                     _str(m.get("message", "")), font_size=10, font_color=P["text"],
                     valign=MSO_ANCHOR.MIDDLE)
    if d.get("call_to_action"):
        _add_textbox(sl, MX, Inches(6.4), MW, Inches(0.4),
                     f"→ {_str(d['call_to_action'])}", font_size=10, font_color=P["gold"], bold=True)
    _slide_footer(sl, P, _extract_refs(d))


# ══════════════════════════════════════════════════════════
# SECTION TYPE → RENDER FUNCTION DISPATCH
# ══════════════════════════════════════════════════════════
SECTION_RENDERERS = {
    "executive_summary": render_executive_summary,
    "prevalence_kpi": render_prevalence_kpi,
    "treatment_algo": render_treatment_algo,
    "guidelines": render_guidelines,
    "unmet_needs": render_unmet_needs,
    "moa": render_moa,
    "pivotal_table": render_pivotal_table,
    "subgroup_analysis": render_subgroup_analysis,
    "competitor_table": render_competitor_table,
    "market_access": render_market_access,
    "swot": render_swot,
    "differentiators": render_differentiators,
    "strategic_imperatives": render_imperatives,
    "imperatives": render_imperatives,
    "narrative": render_narrative,
    "tactical_plan": render_tactics,
    "tactics": render_tactics,
    "areas_interest": render_areas_interest,
    "iep": render_iep,
    "timeline": render_timeline,
    "summary": render_summary,
}


# ══════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════
def _parse_months(raw) -> float:
    s = _str(raw)
    m = re.search(r"([\d.]+)\s*mo", s)
    if m:
        return float(m.group(1))
    m2 = re.search(r"~?([\d.]+)", s)
    if m2:
        try:
            return float(m2.group(1))
        except ValueError:
            pass
    return 0.0


def _theme_key_from_palette(P: dict) -> str:
    for key, pal in PALETTES.items():
        if pal["bg"] == P["bg"]:
            return key
    return "dark"


# ══════════════════════════════════════════════════════════
# MAIN RENDER FUNCTION
# ══════════════════════════════════════════════════════════

# Template slide index → section type mapping (matches the 20-slide templates)
TEMPLATE_SLIDE_MAP = {
    0: "_title",
    1: "executive_summary",
    2: "disease_intro",
    3: "_divider",
    4: "pivotal_table",
    5: "prevalence_kpi",
    6: "differentiators",
    7: "areas_interest",
    8: "iep",
    9: "guidelines",
    10: "moa",
    11: "imperatives",
    12: "treatment_algo",
    13: "tactics",
    14: "narrative",
    15: "unmet_needs",
    16: "timeline",
    17: "market_access",
    18: "swot",
    19: "competitor_table",
}
# Reverse: section_type → slide index
SECTION_TO_SLIDE = {}
for idx, stype in TEMPLATE_SLIDE_MAP.items():
    SECTION_TO_SLIDE[stype] = idx
# Aliases
SECTION_TO_SLIDE["strategic_imperatives"] = 11
SECTION_TO_SLIDE["tactical_plan"] = 13


def _duplicate_slide_internal(pres, slide_index):
    """Duplicate a slide within the same presentation. Safe — no cross-pres issues."""
    src = pres.slides[slide_index]
    layout = src.slide_layout
    new_slide = pres.slides.add_slide(layout)
    
    # Remove auto-generated placeholder shapes
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    
    # Deep copy all shapes from source (same pres → rIds are valid)
    for shape in src.shapes:
        new_slide.shapes._spTree.append(copy_module.deepcopy(shape._element))
    
    # Copy explicit background if set on source
    PML = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    src_bg = src._element.find(f'{PML}cSld/{PML}bg')
    if src_bg is not None:
        dest_cSld = new_slide._element.find(f'{PML}cSld')
        dest_bg = dest_cSld.find(f'{PML}bg')
        if dest_bg is not None:
            dest_cSld.remove(dest_bg)
        dest_cSld.insert(0, copy_module.deepcopy(src_bg))
    
    return len(pres.slides) - 1  # Index of new slide


def _reorder_and_keep(pres, keep_indices):
    """Keep only slides at given indices (in that order), delete the rest."""
    sldIdLst = pres.slides._sldIdLst
    all_items = list(sldIdLst)
    
    # Clear list
    for item in list(sldIdLst):
        sldIdLst.remove(item)
    
    # Re-add in desired order
    for idx in keep_indices:
        sldIdLst.append(all_items[idx])
    
    # Drop relationships for removed slides
    keep_set = set(keep_indices)
    for idx, item in enumerate(all_items):
        if idx not in keep_set:
            rId = item.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', ''
            )
            if rId:
                try:
                    pres.part.drop_rel(rId)
                except Exception:
                    pass


def render_pptx(meta: dict, sections: list[dict], template_id: str = "dark") -> bytes:
    """
    Render a complete premium PPTX from MAP data.

    Two modes:
      1. IN-PLACE: Opens the full 20-slide designer template, fills data, removes
         unused slides, adds extras. No cross-presentation cloning — PowerPoint
         opens without repair warnings.
      2. PROGRAMMATIC: Fallback when no template available. Builds slides from scratch.

    Args:
        meta: dict with drug, indication, year, country, etc.
        sections: list of {id, label, icon, data} dicts from Step 4 review
        template_id: theme key (medai_normal, medai_gold, medai_aquarell, etc.)

    Returns:
        bytes of the generated .pptx file
    """
    # ── Map template_id to internal theme ──
    theme_map = {
        "medai_dark": "normal", "medai_normal": "normal", "normal": "normal",
        "medai_gold": "gold", "gold": "gold",
        "medai_aquarell": "aquarell", "aquarell": "aquarell",
        "dark": "normal", "light": "normal", "gray": "normal",
        "pharma": "normal", "premium": "gold",
        "medai_light": "normal", "consulting": "normal",
        "pharma_blue": "normal", "black_gold": "gold",
    }
    theme = theme_map.get(template_id, "normal")
    P = PALETTES.get(theme, PALETTES["dark"])

    # ── Order sections ──
    exec_item = next((s for s in sections if s.get("id") == "executive_summary"), None)
    summ_item = next((s for s in sections if s.get("id") == "summary"), None)
    others = [s for s in sections if s.get("id") not in ("executive_summary", "summary")]
    ordered = [x for x in [exec_item] + others + [summ_item] if x]

    all_refs = []
    for item in ordered:
        all_refs.extend(_extract_refs(item.get("data", {})))

    # ── Try IN-PLACE mode with full template ──
    full_tpl_path = template_manager.get_full_template_path(theme)
    if full_tpl_path:
        print(f"  ✓ IN-PLACE MODE: {theme}/_full.pptx")
        return _render_inplace(full_tpl_path, meta, ordered, all_refs, P, theme)
    
    # ── Fallback: PROGRAMMATIC mode ──
    print(f"  → PROGRAMMATIC MODE (no _full.pptx for {theme})")
    return _render_programmatic(meta, ordered, all_refs, P, theme)


def _render_inplace(tpl_path, meta, ordered, all_refs, P, theme):
    """
    IN-PLACE rendering: Open full 20-slide template, fill data, reorder, clean up.
    No cross-presentation cloning. PowerPoint opens without warnings.
    """
    # Load template as copy
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    shutil.copy2(str(tpl_path), tmp.name)
    pres = Presentation(tmp.name)
    os.unlink(tmp.name)
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H
    
    DIVIDER_IDX = 3  # Template slide index for divider
    n_original = len(pres.slides)  # Should be 20
    
    # ── Phase 1: Fill title slide (always slide 0) ──
    _fill_title(pres.slides[0], meta, len(ordered), len(set(all_refs)))
    print(f"  ✓ Title slide filled")
    
    # ── Phase 2: Fill template slides that match selected sections ──
    used_template_indices = {0}  # Title is always used
    
    for item in ordered:
        d = item.get("data", {})
        stype = d.get("section_type", item.get("id", ""))
        tpl_idx = SECTION_TO_SLIDE.get(stype)
        
        if tpl_idx is not None and tpl_idx < n_original:
            used_template_indices.add(tpl_idx)
            slide = pres.slides[tpl_idx]
            
            # Try template filler first
            filler = TEMPLATE_FILLERS.get(stype)
            if filler:
                try:
                    filler(slide, d, meta)
                    print(f"  ✓ {stype}: template filled (slide {tpl_idx})")
                except Exception as e:
                    print(f"  ⚠ {stype}: template fill error ({e})")
            else:
                # No specific filler — just set header
                _set_text(_find_shape(slide, "header_title"),
                          item.get("label", stype.replace("_", " ").title()))
                print(f"  ~ {stype}: header-only (slide {tpl_idx})")
    
    # ── Phase 3: Create dividers (duplicate within same pres) ──
    # Each section needs a divider before it. We duplicate the template divider.
    divider_indices = {}  # section_type → new slide index
    
    for sec_num, item in enumerate(ordered, 1):
        d = item.get("data", {})
        stype = d.get("section_type", item.get("id", ""))
        label = item.get("label", stype.replace("_", " ").title())
        
        new_idx = _duplicate_slide_internal(pres, DIVIDER_IDX)
        _fill_divider(pres.slides[new_idx], sec_num, label)
        divider_indices[stype] = new_idx
    
    print(f"  ✓ {len(divider_indices)} dividers created")
    
    # ── Phase 4: Add extra programmatic slides (charts, overflow tables) ──
    # These are sections that produce ADDITIONAL slides beyond the template slide.
    # They append to the presentation and we track their indices.
    extra_indices = {}  # section_type → [list of extra slide indices]
    
    for item in ordered:
        d = item.get("data", {})
        stype = d.get("section_type", item.get("id", ""))
        extras = []
        theme_key = _theme_key_from_palette(P)
        
        try:
            if stype == "pivotal_table":
                # Extra slides for KM curves (per trial beyond the first)
                trials = d.get("trials", [])
                active = [t for t in trials if t.get("include_in_map") is not False]
                if len(active) > 1:
                    # Template slide handles first trial; render extra trials programmatically
                    overflow_start = len(pres.slides)
                    # Skip first trial (already in template), render rest
                    for tr in active[1:]:
                        try:
                            render_pivotal_table(pres, {"trials": [tr], "section_type": "pivotal_table"}, P, meta)
                        except Exception:
                            pass
                    extras = list(range(overflow_start, len(pres.slides)))
            
            elif stype == "competitor_table":
                # Extra slides: mPFS chart + ORR chart + table pages
                overflow_start = len(pres.slides)
                render_competitor_table(pres, d, P, meta)
                extras = list(range(overflow_start, len(pres.slides)))
                # Keep template slide too (has header + source filled by filler)
            
            elif stype == "subgroup_analysis":
                overflow_start = len(pres.slides)
                render_subgroup_analysis(pres, d, P, meta)
                extras = list(range(overflow_start, len(pres.slides)))
            
            elif stype == "summary":
                overflow_start = len(pres.slides)
                render_summary(pres, d, P, meta)
                extras = list(range(overflow_start, len(pres.slides)))
        
        except Exception as e:
            print(f"  ⚠ {stype}: extra slides error ({e})")
        
        if extras:
            extra_indices[stype] = extras
            print(f"  + {stype}: {len(extras)} extra slides")
    
    # ── Phase 5: Build final slide order ──
    final_order = [0]  # Title slide
    
    for item in ordered:
        d = item.get("data", {})
        stype = d.get("section_type", item.get("id", ""))
        
        # Add divider
        div_idx = divider_indices.get(stype)
        if div_idx is not None:
            final_order.append(div_idx)
        
        # Add content slide(s)
        if stype in extra_indices:
            # Use programmatic slides instead of template slide
            final_order.extend(extra_indices[stype])
        else:
            # Use template slide
            tpl_idx = SECTION_TO_SLIDE.get(stype)
            if tpl_idx is not None and tpl_idx in used_template_indices:
                final_order.append(tpl_idx)
            elif tpl_idx is None and stype not in extra_indices:
                # Unknown section — try programmatic fallback
                renderer = SECTION_RENDERERS.get(stype)
                if renderer:
                    try:
                        fallback_start = len(pres.slides)
                        renderer(pres, d, P, meta)
                        for fi in range(fallback_start, len(pres.slides)):
                            final_order.append(fi)
                    except Exception as e:
                        print(f"  ⚠ {stype}: fallback render error ({e})")
    
    # ── Phase 6: Add references slide ──
    ref_layout = pres.slides[0].slide_layout  # Use same layout as title for master inheritance
    ref_sl = pres.slides.add_slide(ref_layout)
    # Clear placeholder shapes
    for shape in list(ref_sl.shapes):
        shape._element.getparent().remove(shape._element)
    
    _slide_header(ref_sl, "References", P, f"MedAI Suite · {len(set(all_refs))} sources")
    unique_refs = list(dict.fromkeys(all_refs))[:30]
    ref_text = "\n".join(f"[{i+1}] {r}" for i, r in enumerate(unique_refs) if r)
    _add_textbox(ref_sl, MX, Inches(1.1), MW, Inches(5.8),
                 ref_text or "No references collected.",
                 font_size=8, font_color=P["muted"], font_name="DM Mono")
    final_order.append(len(pres.slides) - 1)
    
    # ── Phase 7: Reorder and remove unused slides ──
    print(f"  Final: {len(final_order)} slides (from {len(pres.slides)} total)")
    _reorder_and_keep(pres, final_order)
    
    # Write to bytes
    buf = io.BytesIO()
    pres.save(buf)
    buf.seek(0)
    print(f"  ✅ IN-PLACE render complete: {buf.getbuffer().nbytes // 1024} KB")
    return buf.read()


def _render_programmatic(meta, ordered, all_refs, P, theme):
    """
    PROGRAMMATIC rendering: Fallback when no full template available.
    Builds all slides from scratch. Works but no designer backgrounds.
    """
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H

    # ── Title slide ──
    title_sl = _new_slide(pres, P)
    _add_rect(title_sl, Inches(0), Inches(0), SLIDE_W, SLIDE_H, P["bg"])
    _add_rect(title_sl, MX, Inches(2.2), Inches(4), Inches(0.06), P["teal"])
    _add_textbox(title_sl, MX, Inches(2.5), MW, Inches(1.2),
                 meta.get("drug", "Medical Affairs Plan"),
                 font_size=36, font_color=P["white"], bold=True)
    scope = meta.get("country") or meta.get("region") or "Global"
    _add_textbox(title_sl, MX, Inches(3.6), MW, Inches(0.6),
                 f"{scope} Medical Affairs Plan — {meta.get('year', '')}",
                 font_size=16, font_color=P["muted"])
    _add_textbox(title_sl, MX, Inches(4.3), MW, Inches(0.4),
                 meta.get("indication", ""), font_size=12, font_color=P["teal"])
    _add_textbox(title_sl, MX, Inches(5.2), MW, Inches(0.4),
                 f"{len(ordered)} Sections · {len(set(all_refs))} References · {meta.get('model', 'Claude')}",
                 font_size=10, font_color=P["dim"], font_name="DM Mono")
    _add_textbox(title_sl, MX, Inches(6.8), MW, Inches(0.3),
                 "MedAI Suite Premium · AI-Verified by ELISE",
                 font_size=8, font_color=P["dim"], font_name="DM Mono")

    # ── Content slides ──
    for sec_num, item in enumerate(ordered, 1):
        d = item.get("data", {})
        stype = d.get("section_type", item.get("id", ""))
        label = item.get("label", stype)

        # Divider
        _add_divider_slide(pres, sec_num, label, P)

        # Content
        renderer = SECTION_RENDERERS.get(stype)
        if renderer:
            try:
                renderer(pres, d, P, meta)
            except Exception as e:
                print(f"  ✗ {label}: render error ({e})")
                render_generic(pres, d, P, meta, label)
        else:
            render_generic(pres, d, P, meta, label)

    # ── References ──
    ref_sl = _new_slide(pres, P)
    ref_sl.background.fill.solid()
    ref_sl.background.fill.fore_color.rgb = _rgb(P["bg"])
    _slide_header(ref_sl, "References", P, f"MedAI Suite · {len(set(all_refs))} sources")
    unique_refs = list(dict.fromkeys(all_refs))[:30]
    ref_text = "\n".join(f"[{i+1}] {r}" for i, r in enumerate(unique_refs) if r)
    _add_textbox(ref_sl, MX, Inches(1.1), MW, Inches(5.8),
                 ref_text or "No references collected.",
                 font_size=8, font_color=P["muted"], font_name="DM Mono")

    buf = io.BytesIO()
    pres.save(buf)
    buf.seek(0)
    print(f"  ✅ PROGRAMMATIC render complete: {buf.getbuffer().nbytes // 1024} KB")
    return buf.read()
